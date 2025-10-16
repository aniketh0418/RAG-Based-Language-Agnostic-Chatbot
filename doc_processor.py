import os
import io
import uuid
import logging
from typing import List, Tuple, Dict

import pymupdf as fitz
from PIL import Image
import docx
from pptx import Presentation
from fpdf import FPDF
import tempfile
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.http import models
from google.generativeai import configure, GenerativeModel

# ---------------- CONFIG ----------------
QDRANT_API_KEY = "your-qdrant-api-key"
QDRANT_URL = "your-qdrant-connection-url"
QDRANT_COLLECTION = "your-qdrant-collection-name"
EMBEDDING_MODEL_NAME = "BAAI/bge-small-en-v1.5"

PAGE_TEXT_CHAR_THRESHOLD = 30
IMAGE_DPI_SCALE = 2

# Gemini OCR setup
configure(api_key="your-gemini-api-key") 
gemini_model = GenerativeModel("gemini-2.5-flash")

PROMPT = """
You are given an image or document that contains tabular data (such as a timetable, exam schedule, financial table, dataset, or any structured information).  

Your task:
1. Extract all data from the table while preserving the structure and semantic meaning.  
2. Present the output only in JSON format.  
3. Ensure the JSON is well-structured for retrieval-augmented generation (RAG) use.  
4. Each row of the table should be a JSON object with column headers as keys.  
5. Include all details exactly as written (dates, times, codes, abbreviations).  
6. Output valid JSON only.
"""

# --------- logging ---------
logging.basicConfig(level=logging.INFO)
log = logging.getLogger("ingest")

# --------- clients & models ---------
client = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY, timeout=60)
embedding_model = SentenceTransformer(EMBEDDING_MODEL_NAME)


# --------- Qdrant helper ---------
def ensure_collection(collection_name: str, vector_size: int):
    try:
        client.get_collection(collection_name=collection_name)
    except Exception:
        client.create_collection(
            collection_name=collection_name,
            vectors_config=models.VectorParams(size=vector_size, distance=models.Distance.COSINE),
        )


# --------- Gemini OCR ---------
def extract_text_with_gemini(image_path: str) -> str:
    mime = "image/png"
    if image_path.lower().endswith((".jpg", ".jpeg")):
        mime = "image/jpeg"

    image_file = {"mime_type": mime, "data": open(image_path, "rb").read()}
    prompt = "Extract all readable text from the provided image."
    response = gemini_model.generate_content([prompt, image_file])
    return response.text.strip() if response and response.text else ""


# --------- Timetable extractor ---------
def extract_timetable(file_path: str) -> List[str]:
    ext = os.path.splitext(file_path)[-1].lower()
    results = []

    if ext == ".pdf":
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, start=1):
            pix = page.get_pixmap(dpi=200)
            img_data = pix.tobytes("png")
            image_file = {"mime_type": "image/png", "data": img_data}
            response = gemini_model.generate_content([PROMPT, image_file])
            extracted_text = response.text.strip() if response.text else "[No text extracted]"
            results.append(extracted_text)

    elif ext in [".png", ".jpg", ".jpeg"]:
        with open(file_path, "rb") as f:
            img_data = f.read()
        image_file = {"mime_type": f"image/{ext[1:]}", "data": img_data}
        response = gemini_model.generate_content([PROMPT, image_file])
        extracted_text = response.text.strip() if response.text else "[No text extracted]"
        results.append(extracted_text)

    else:
        raise ValueError("Unsupported file format for timetable extraction")

    return results


# --------- Extractors ---------
def extract_pdf_pages_text_or_image(path: str) -> List[Tuple[int, str, bool]]:
    doc = fitz.open(path)
    pages_out = []
    for i, page in enumerate(doc):
        raw_text = page.get_text("text").strip()
        if len(raw_text) < PAGE_TEXT_CHAR_THRESHOLD:
            mat = fitz.Matrix(IMAGE_DPI_SCALE, IMAGE_DPI_SCALE)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            tmp_img = tempfile.mktemp(suffix=".png")
            pix.save(tmp_img)
            text = extract_text_with_gemini(tmp_img)
            pages_out.append((i + 1, text, True))
        else:
            pages_out.append((i + 1, raw_text, False))
    return pages_out


def extract_image_text(path: str) -> List[Tuple[int, str, bool]]:
    text = extract_text_with_gemini(path)
    return [(1, text, True)]


def convert_docx_to_pdf(path: str) -> str:
    doc = docx.Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for para in paras:
        pdf.multi_cell(0, 10, para)
    tmp_pdf = tempfile.mktemp(suffix=".pdf")
    pdf.output(tmp_pdf)
    return tmp_pdf


def convert_pptx_to_pdf(path: str) -> str:
    prs = Presentation(path)
    pdf = FPDF()
    pdf.set_font("Arial", size=12)
    for slide in prs.slides:
        pdf.add_page()
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        pdf.multi_cell(0, 10, "\n".join(texts))
    tmp_pdf = tempfile.mktemp(suffix=".pdf")
    pdf.output(tmp_pdf)
    return tmp_pdf


# --------- Wrapper ---------
def extract_text_from_file(path: str, table_data: bool = False) -> List[Tuple[int, str, bool]]:
    if table_data:
        timetable_texts = extract_timetable(path)
        return [(i + 1, t, True) for i, t in enumerate(timetable_texts)]

    lower = path.lower()
    if lower.endswith(".pdf"):
        return extract_pdf_pages_text_or_image(path)
    elif lower.endswith(".docx"):
        pdf_path = convert_docx_to_pdf(path)
        return extract_pdf_pages_text_or_image(pdf_path)
    elif lower.endswith(".pptx"):
        pdf_path = convert_pptx_to_pdf(path)
        return extract_pdf_pages_text_or_image(pdf_path)
    elif lower.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp")):
        return extract_image_text(path)
    else:
        raise ValueError("Unsupported file type: " + path)


# --------- Chunking ---------
def chunk_pages(pages: List[Tuple[int, str, bool]], chunk_size=800, chunk_overlap=150):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    all_chunks = []
    for (page_num, text, used_ocr) in pages:
        if not text or not text.strip():
            continue
        parts = splitter.split_text(text)
        for idx, chunk in enumerate(parts):
            if len(chunk.strip()) < 20:
                continue
            metadata = {
                "page": page_num,
                "chunk_index": idx,
                "chunk_length": len(chunk),
                "ocr_used": used_ocr,
            }
            all_chunks.append((chunk, metadata))
    return all_chunks


# --------- Qdrant upsert ---------
def store_chunks_in_qdrant(chunks: List[Tuple[str, dict]], collection_name: str,
                           doc_id: str, file_name: str, batch_size=64):
    texts = [c[0] for c in chunks]
    metas = [c[1] for c in chunks]
    if not texts:
        return 0

    vectors = embedding_model.encode(texts, show_progress_bar=True)
    vec_dim = len(vectors[0])
    ensure_collection(collection_name, vec_dim)

    total = 0
    for i in range(0, len(texts), batch_size):
        slice_texts = texts[i:i + batch_size]
        slice_metas = metas[i:i + batch_size]
        slice_vectors = vectors[i:i + batch_size]

        points = []
        for t, m, v in zip(slice_texts, slice_metas, slice_vectors):
            m["doc_id"] = doc_id
            m["file_name"] = file_name  # ✅ consistent for Qdrant filtering
            payload = {"text": t, "metadata": m}
            point = models.PointStruct(
                id=uuid.uuid4().int >> 64,
                vector=v.tolist() if hasattr(v, "tolist") else list(v),
                payload=payload
            )
            points.append(point)

        client.upsert(collection_name=collection_name, points=points)
        total += len(points)

    return total


def process_file_to_qdrant(file_path: str, collection_name=QDRANT_COLLECTION,
                           doc_id: str = None, table_data: bool = False, file_name: str = None) -> Dict:
    if doc_id is None:
        doc_id = str(uuid.uuid4())
    if file_name is None:
        file_name = os.path.basename(file_path)

    log.info(f"Processing file: {file_path}")
    pages = extract_text_from_file(file_path, table_data=table_data)
    log.info(f"Extracted {len(pages)} page items.")

    chunks = chunk_pages(pages)
    log.info(f"Created {len(chunks)} chunks from file.")

    stored = store_chunks_in_qdrant(chunks, collection_name, doc_id, file_name)
    log.info(f"Stored {stored} chunks into collection '{collection_name}' for doc_id={doc_id}.")

    return {
        "doc_id": doc_id,
        "chunks_stored": stored,
        "collection_name": collection_name,
        "source_file": file_name,
    }


# --------- CLI ---------
if __name__ == "__main__":
    file_to_process = "tt_pdf.pdf"
    collection_name = QDRANT_COLLECTION
    doc_id = None
    result = process_file_to_qdrant(
        file_to_process,
        collection_name=collection_name,
        doc_id=doc_id,
        table_data=True
    )

    print("RESULT:", result)

